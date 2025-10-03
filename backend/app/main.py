# backend/app/main.py
# REAL AI-powered backend with Gemini - COMPLETE FIXED VERSION

from fastapi import FastAPI, HTTPException, UploadFile, File
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse, RedirectResponse
from pydantic import BaseModel
from typing import List, Dict, Any, Optional
import asyncio
import uuid
from datetime import datetime
import json
import io
import base64
import re
import os
import random
from pathlib import Path


try:
    from dotenv import load_dotenv
    load_dotenv()
except ImportError:
    print("⚠️ python-dotenv not installed. Run: pip install python-dotenv")
    
# Document processing
try:
    import PyPDF2
    HAS_PDF = True
except ImportError:
    HAS_PDF = False
    print("⚠️ PyPDF2 not installed: pip install PyPDF2")

try:
    from docx import Document
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False
    print("⚠️ python-docx not installed: pip install python-docx")

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, RGBColor
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False
    print("⚠️ python-pptx not installed: pip install python-pptx")

# AI Models
AVAILABLE_AI = []

# Google Gemini (PRIORITIZED)
try:
    import google.generativeai as genai
    # Use your new working API key
    api_key = os.getenv("GOOGLE_API_KEY", "AIzaSyAs2JG53pWD7mbFwoOss8kVEkbvimw8Jyw")
    if api_key:
        genai.configure(api_key=api_key)
        AVAILABLE_AI.append("gemini")
        print("✅ Google Gemini API available")
except ImportError:
    print("⚠️ Google Gemini not installed: pip install google-generativeai")
except Exception as e:
    print(f"⚠️ Gemini configuration error: {e}")

# OpenAI
try:
    import openai
    from openai import OpenAI
    if os.getenv("OPENAI_API_KEY"):
        AVAILABLE_AI.append("openai")
        print("✅ OpenAI API available")
except ImportError:
    pass

# Transformers
try:
    from transformers import pipeline
    AVAILABLE_AI.append("transformers")
    print("✅ Local Transformers available")
except ImportError:
    pass

app = FastAPI(
    title="Prompt-to-Pipeline API - Gemini Enhanced",
    description="AI-powered document processing with Google Gemini",
    version="4.0.0"
)

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Storage
uploaded_documents = {}
pipelines = {}

# ================== AI Service ==================
class RealAIService:
    def __init__(self):
        self.method = self._detect_best_method()
        self._initialize_models()
        
    def _detect_best_method(self):
        if "gemini" in AVAILABLE_AI:
            return "gemini"
        elif "openai" in AVAILABLE_AI:
            return "openai"
        elif "transformers" in AVAILABLE_AI:
            return "transformers"
        return "basic"
    
    def _initialize_models(self):
        if self.method == "gemini":
            self.gemini_api_key = os.getenv("GOOGLE_API_KEY")
            if not api_key:
                print("⚠️ Warning: GOOGLE_API_KEY not set in environment")
            
            # Try REST API with the actual available models
            try:
                import requests
                
                # Try models in order of preference (fastest/best first)
                model_options = [
                    'gemini-2.5-flash',      # Newest, fastest
                    'gemini-2.5-flash-lite', # Fastest, lightweight
                    'gemini-2.0-flash',      # Fallback
                ]
                
                for model_name in model_options:
                    try:
                        test_url = f"https://generativelanguage.googleapis.com/v1/models/{model_name}:generateContent?key={self.gemini_api_key}"
                        test_data = {
                            "contents": [{
                                "parts": [{"text": "Hello"}]
                            }]
                        }
                        
                        response = requests.post(test_url, json=test_data, timeout=10)
                        
                        if response.status_code == 200:
                            print(f"✅ Gemini REST API working with {model_name}")
                            self.use_rest_api = True
                            self.gemini_model = model_name
                            return
                        else:
                            print(f"❌ {model_name} returned: {response.status_code}")
                    except Exception as e:
                        print(f"❌ {model_name} error: {str(e)[:80]}")
                        continue
                    
            except Exception as e:
                print(f"❌ Gemini REST API setup error: {e}")
            
            print("⚠️ Gemini API not working, using enhanced fallback")
            self.method = "basic"
            self.use_rest_api = False
                
        elif self.method == "openai":
            self.client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))
    
    def _call_gemini_rest(self, prompt: str, max_tokens: int = 1000, temperature: float = 0.4) -> str:
        """Call Gemini using REST API - handles Gemini 2.5"""
        import requests
        
        model_name = getattr(self, 'gemini_model', 'gemini-2.5-flash')
        url = f"https://generativelanguage.googleapis.com/v1/models/{model_name}:generateContent?key={self.gemini_api_key}"
        
        # For Gemini 2.5, we need MUCH higher token limits because of thinking mode
        # The model uses tokens for internal reasoning before generating output
        if '2.5' in model_name:
            # Double the requested tokens to account for thinking
            max_tokens = max_tokens * 2
        
        data = {
            "contents": [{
                "parts": [{"text": prompt}]
            }],
            "generationConfig": {
                "temperature": temperature,
                "maxOutputTokens": max_tokens,
                "topP": 0.8,
                "topK": 40
            }
        }
        
        response = requests.post(url, json=data, timeout=40)
        
        if response.status_code == 200:
            result = response.json()
            
            if 'candidates' in result and len(result['candidates']) > 0:
                candidate = result['candidates'][0]
                
                # Get the actual text content
                if 'content' in candidate and 'parts' in candidate['content']:
                    parts = candidate['content']['parts']
                    if len(parts) > 0 and 'text' in parts[0]:
                        text = parts[0]['text']
                        
                        # Check if truncated
                        finish_reason = candidate.get('finishReason', '')
                        if finish_reason == 'MAX_TOKENS':
                            print(f"  ⚠️ Response truncated at {max_tokens} tokens")
                        
                        return text
            
            # Response was empty or malformed
            print(f"  ❌ Empty Gemini response")
            raise Exception("Empty response from Gemini API")
        else:
            raise Exception(f"Gemini API {response.status_code}: {response.text}")
    
    async def summarize(self, text: str, length: str = "medium") -> str:
        max_chars = 20000
        if len(text) > max_chars:
            text = text[:max_chars]
        
        try:
            if self.method == "gemini":
                prompt = f"""You are an expert at creating comprehensive, informative summaries.

Create a well-structured summary of this document that captures:
- The main topic and purpose
- Key findings, results, or conclusions
- Important methodologies or approaches mentioned
- Specific details, numbers, or facts
- Overall significance or implications

Length: {length} detail
- short: 2-3 focused paragraphs
- medium: 4-6 detailed paragraphs with specific information
- long: 7-10 comprehensive paragraphs covering all aspects

Make the summary informative and useful. Include specific details from the document.

Document:
{text}

Write a {length}-length summary:"""
                
                if hasattr(self, 'use_rest_api') and self.use_rest_api:
                    summary = self._call_gemini_rest(prompt, max_tokens=1500)
                else:
                    response = self.model.generate_content(
                        prompt,
                        generation_config=genai.types.GenerationConfig(
                            temperature=0.4,
                            top_p=0.85,
                            max_output_tokens=1500
                        )
                    )
                    summary = response.text.strip()
                
                if len(summary) < 200:
                    print("⚠️ Summary too short, using enhanced extraction")
                    return self._enhanced_summary(text)
                
                return summary
            
            return self._enhanced_summary(text)
        except Exception as e:
            print(f"Summarization error: {e}")
            import traceback
            traceback.print_exc()
            return self._enhanced_summary(text)
    
    def _enhanced_summary(self, text: str) -> str:
        """MUCH BETTER extraction-based summary - plain text, no formatting"""
        print("  📊 Creating intelligent extraction-based summary...")
        
        sentences = [s.strip() + '.' for s in text.split('.') if len(s.strip()) > 30]
        
        if len(sentences) < 5:
            return "Document too short to summarize effectively."
        
        summary_parts = []
        
        # Part 1: Introduction (first 2-3 sentences)
        intro = ' '.join(sentences[:3])
        summary_parts.append(f"Introduction: {intro}")
        
        # Part 2: Key Findings and Results
        result_keywords = ['result', 'found', 'showed', 'demonstrated', 'achieved', 
                          'accuracy', 'performance', 'concluded', 'reveals', 'indicates',
                          'obtained', 'reached', 'attained']
        
        result_sents = []
        for s in sentences:
            if any(kw in s.lower() for kw in result_keywords):
                result_sents.append(s)
                if len(result_sents) >= 4:
                    break
        
        if result_sents:
            summary_parts.append(f"Key Findings: {' '.join(result_sents)}")
        
        # Part 3: Methodology
        method_keywords = ['method', 'approach', 'technique', 'algorithm', 'model',
                          'used', 'applied', 'implemented', 'developed', 'proposed',
                          'dataset', 'trained', 'evaluated']
        
        method_sents = []
        for s in sentences:
            if any(kw in s.lower() for kw in method_keywords) and s not in result_sents:
                method_sents.append(s)
                if len(method_sents) >= 3:
                    break
        
        if method_sents:
            summary_parts.append(f"Methodology: {' '.join(method_sents)}")
        
        # Part 4: Conclusion (last 2 sentences)
        conclusion = ' '.join(sentences[-2:])
        summary_parts.append(f"Conclusion: {conclusion}")
        
        return '\n\n'.join(summary_parts)
    
    async def extract_key_points(self, text: str) -> List[str]:
        try:
            if self.method == "gemini":
                prompt = f"""Extract 8-12 key points from this document.

Requirements for each point:
- Must be a complete, specific insight or finding
- Include concrete details (numbers, names, methods, results)
- Focus on main concepts, findings, conclusions, and methodologies
- Be informative and standalone (someone reading just the point should learn something specific)
- Avoid vague statements like "The document discusses X"
- Each point: 1-3 sentences maximum

Document:
{text[:12000]}

List the key points (format: • Point text):"""
                
                # Use REST API if available
                if hasattr(self, 'use_rest_api') and self.use_rest_api:
                    response_text = self._call_gemini_rest(prompt, max_tokens=2000, temperature=0.2)
                else:
                    response = self.model.generate_content(
                        prompt,
                        generation_config=genai.types.GenerationConfig(
                            temperature=0.3,
                            max_output_tokens=800
                        )
                    )
                    response_text = response.text
                
                points = []
                for line in response_text.split('\n'):
                    cleaned = line.strip('- •*0123456789. \t')
                    # Filter out section headers and short points
                    if (cleaned and 
                        len(cleaned) > 30 and 
                        not cleaned.endswith(':') and
                        not cleaned.isupper()):
                        points.append(cleaned)
                
                if len(points) >= 5:
                    print(f"  ✅ Extracted {len(points)} quality points from Gemini")
                    return points[:12]
                else:
                    print(f"  ⚠️ Only {len(points)} points from Gemini, using enhanced extraction")
            
            return self._enhanced_key_points(text)
        except Exception as e:
            print(f"Key points error: {e}")
            import traceback
            traceback.print_exc()
            return self._enhanced_key_points(text)
    
    def _enhanced_key_points(self, text: str) -> List[str]:
        """Enhanced keyword-based extraction"""
        sentences = [s.strip() for s in text.split('.') if len(s.strip()) > 40]
        points = []
        
        # Priority 1: Sentences with results/findings
        high_priority = ['result', 'found', 'showed', 'demonstrated', 'achieved', 
                        'accuracy', 'performance', 'concluded', 'indicates', 'reveals']
        
        for s in sentences:
            if any(w in s.lower() for w in high_priority) and 60 < len(s) < 250:
                if s not in points:
                    points.append(s)
                    if len(points) >= 12:
                        break
        
        # Priority 2: Sentences with methodology
        method_keywords = ['method', 'approach', 'technique', 'algorithm', 'model', 
                          'system', 'framework', 'architecture', 'used', 'applied']
        
        for s in sentences:
            if any(w in s.lower() for w in method_keywords) and 60 < len(s) < 250:
                if s not in points and len(points) < 12:
                    points.append(s)
        
        # Priority 3: Sentences with importance markers
        important = ['important', 'significant', 'key', 'main', 'critical', 'essential',
                    'primary', 'major', 'fundamental']
        
        for s in sentences:
            if any(w in s.lower() for w in important) and 60 < len(s) < 250:
                if s not in points and len(points) < 12:
                    points.append(s)
        
        # If still not enough, add well-formed sentences
        if len(points) < 8:
            for s in sentences:
                if 70 < len(s) < 200 and s not in points:
                    points.append(s)
                    if len(points) >= 10:
                        break
        
        return points[:12]
    
    async def generate_quiz(self, text: str, num_questions: int = 5) -> List[Dict]:
        print(f"\n🎯 Generating {num_questions} quiz questions...")
        
        # Try Gemini first if available
        if self.method == "gemini" and hasattr(self, 'use_rest_api') and self.use_rest_api:
            try:
                prompt = f"""Create {num_questions} multiple choice quiz questions from this document.

RULES:
1. Questions must test SPECIFIC content from the document
2. Each question has 4 options
3. Only ONE correct answer
4. Wrong answers must be plausible but clearly incorrect
5. NO duplicate questions
6. Include variety: definitions, numbers, concepts, applications

Document:
{text[:10000]}

Return ONLY a JSON array, no other text:
[
  {{"question": "What is X according to the document?", "options": ["correct", "wrong1", "wrong2", "wrong3"], "correct_answer": 0, "explanation": "..."}}
]"""
                
                print("  📡 Calling Gemini for quiz generation...")
                response_text = self._call_gemini_rest(prompt, max_tokens=5000, temperature=0.8)
                print(f"  ✅ Gemini response received ({len(response_text)} chars)")
                
                # Extract JSON
                json_match = re.search(r'\[[\s\S]*\]', response_text)
                if json_match:
                    quiz_data = json.loads(json_match.group())
                    
                    if len(quiz_data) >= num_questions:
                        print(f"  ✅ Gemini generated {len(quiz_data)} questions successfully")
                        
                        # Validate and return
                        valid_quiz = []
                        for q in quiz_data[:num_questions]:
                            if isinstance(q, dict) and 'question' in q and 'options' in q:
                                if 'correct_answer' not in q:
                                    q['correct_answer'] = 0
                                elif isinstance(q['correct_answer'], str):
                                    q['correct_answer'] = ord(q['correct_answer'].upper()) - ord('A')
                                
                                while len(q['options']) < 4:
                                    q['options'].append(f"Option {len(q['options'])+1}")
                                q['options'] = q['options'][:4]
                                
                                if 'explanation' not in q:
                                    q['explanation'] = "Based on document."
                                
                                valid_quiz.append(q)
                        
                        if valid_quiz:
                            return valid_quiz
                    else:
                        print(f"  ⚠️ Only got {len(quiz_data)} questions from Gemini")
                else:
                    print("  ⚠️ No JSON found in Gemini response")
                    print(f"  Response preview: {response_text[:300]}")
            
            except Exception as e:
                print(f"  ❌ Gemini quiz error: {str(e)[:150]}")
                import traceback
                traceback.print_exc()
        
        # Use smart fallback
        print("  🔄 Using intelligent context-based quiz generation...")
        return self._smart_context_quiz(text, num_questions)
    
    def _smart_context_quiz(self, text: str, num: int) -> List[Dict]:
        """Generate intelligent quiz questions based on document analysis"""
        print(f"🧠 Smart context quiz generation for {num} questions...")
        questions = []
        
        sentences = [s.strip() for s in text.split('.') if 50 < len(s.strip()) < 300]
        
        # Strategy 1: Extract REAL definitions with MUCH better parsing
        # Pattern 1: "ACRONYM stands for Full Name"
        acronym_pattern = r'\b([A-Z]{2,8})\s+stands for\s+([A-Z][^.]{10,120})\.'
        acronym_matches = re.findall(acronym_pattern, text)
        
        # Pattern 2: "Term is a/an definition"
        definition_pattern = r'\b([A-Z][a-z]{2,}(?:\s+[A-Z][a-z]+){0,3})\s+is\s+(?:a|an)\s+([a-z][^.]{20,150})\.'
        definition_matches = re.findall(definition_pattern, text)
        
        # Pattern 3: "Term refers to/means definition"
        refers_pattern = r'\b([A-Z][a-z]{2,}(?:\s+[A-Z][a-z]+){0,3})\s+(?:refers to|means)\s+([^.]{15,120})\.'
        refers_matches = re.findall(refers_pattern, text)
        
        # Combine all definitions
        all_definitions = []
        
        # Process acronyms
        for acronym, full_name in acronym_matches:
            if len(acronym) >= 2 and len(full_name) > 10:
                all_definitions.append((acronym, f"stands for {full_name}", True))
        
        # Process regular definitions
        for subject, definition in definition_matches + refers_matches:
            # Validate subject
            if (len(subject) > 2 and 
                subject not in ['WHAT', 'WHO', 'WHERE', 'WHEN', 'WHY', 'HOW', 'This', 'That', 'These', 'Those'] and
                'page' not in subject.lower() and
                'figure' not in subject.lower() and
                len(definition) > 15):
                all_definitions.append((subject.strip(), definition.strip(), False))
        
        print(f"  Found {len(all_definitions)} validated definitions")
        
        # Create questions from definitions
        for subject, definition, is_acronym in all_definitions[:min(num, len(all_definitions))]:
            
            # Generate smart wrong answers
            wrong_answers = []
            
            # Use other definitions from document
            other_defs = [d for s, d, _ in all_definitions if s != subject][:2]
            wrong_answers.extend(other_defs)
            
            # Add contextual wrong answers based on document
            if is_acronym:
                # For acronyms, create plausible expansions
                words = definition.replace('stands for ', '').split()
                if len(words) >= 3:
                    # Shuffle words for wrong answer
                    shuffled = words.copy()
                    random.shuffle(shuffled)
                    wrong_answers.append('stands for ' + ' '.join(shuffled[:len(words)]))
                
                wrong_answers.extend([
                    'stands for a General Processing Technique',
                    'stands for a type of machine learning model'
                ])
            else:
                # For regular terms, use related concepts
                tech_phrases = re.findall(r'(?:a|an)\s+([a-z][^.,]{15,60}(?:model|system|method|technique))', text, re.IGNORECASE)
                for phrase in tech_phrases[:3]:
                    cleaned = phrase.strip()
                    if cleaned != definition and cleaned not in wrong_answers:
                        wrong_answers.append(cleaned)
                
                # Generic but plausible
                wrong_answers.extend([
                    'a statistical analysis method',
                    'a data preprocessing technique',
                    'a type of neural network'
                ])
            
            # Ensure we have exactly 3 wrong answers that are different from correct
            final_wrong = []
            for wa in wrong_answers:
                if wa != definition and wa not in final_wrong and len(final_wrong) < 3:
                    final_wrong.append(wa)
            
            while len(final_wrong) < 3:
                final_wrong.append(f"Option {len(final_wrong) + 1}")
            
            # Create options
            if is_acronym:
                correct_answer = definition
            else:
                correct_answer = definition
            
            options = [correct_answer] + final_wrong[:3]
            random.shuffle(options)
            correct_idx = options.index(correct_answer)
            
            # Create question
            if is_acronym:
                question_text = f"What does {subject} stand for?"
            else:
                question_text = f"What is {subject}?"
            
            questions.append({
                "question": question_text,
                "options": options,
                "correct_answer": correct_idx,
                "explanation": f"{subject} {definition}"
            })
            
            print(f"  ✓ Q: {question_text}")
        
        # Strategy 2: Numerical facts with context
        # Find sentences with numbers and meaningful context
        number_pattern = r'([^.]*?(\d+(?:,\d+)*(?:\.\d+)?)\s*(?:percent|%|years?|months?|participants?|cases?|times?|million|billion|thousand|accuracy|precision|recall)[^.]*?\.)'
        number_facts = re.findall(number_pattern, text)
        
        print(f"  Found {len(number_facts)} numerical facts")
        
        for fact_text, number in number_facts[:num - len(questions)]:
            fact_text = fact_text.strip()
            
            # Create a specific question about this fact
            # Remove the number to create the question
            question_base = fact_text.replace(number, "____")
            
            # Generate realistic wrong numbers
            try:
                if '.' in number or ',' in number:
                    base = float(number.replace(',', ''))
                    if base > 0.5 and base < 1:  # Percentage as decimal
                        wrong_nums = [
                            f"{base + 0.15:.2f}",
                            f"{base - 0.10:.2f}",
                            f"{base * 1.3:.2f}"
                        ]
                    else:
                        wrong_nums = [
                            f"{base * 1.25:.1f}",
                            f"{base * 0.85:.1f}",
                            f"{base * 1.5:.1f}"
                        ]
                else:
                    base = int(number.replace(',', ''))
                    if base < 100:
                        wrong_nums = [str(base + 15), str(base - 10), str(base + 25)]
                    else:
                        wrong_nums = [str(int(base * 1.2)), str(int(base * 0.8)), str(int(base * 1.5))]
                
                # Remove duplicates and ensure they're different from correct
                wrong_nums = [w for w in wrong_nums if w != number][:3]
                
            except:
                wrong_nums = ["42", "75", "128"]
            
            options = [number] + wrong_nums
            random.shuffle(options)
            correct_idx = options.index(number)
            
            # Create a natural question
            if "percent" in fact_text.lower() or "%" in fact_text:
                q_text = f"What percentage is mentioned in: \"{question_base}\"?"
            elif "year" in fact_text.lower():
                q_text = f"How many years are mentioned: \"{question_base}\"?"
            else:
                q_text = f"What is the correct number: \"{question_base}\"?"
            
            questions.append({
                "question": q_text,
                "options": options,
                "correct_answer": correct_idx,
                "explanation": f"The document states: {fact_text}"
            })
            
            print(f"  ✓ Q: Number question ({number})")
        
        # Strategy 3: Find key concepts and their descriptions
        # Look for important technical terms or concepts
        concept_pattern = r'\b([A-Z][a-z]+(?:\s+[A-Z][a-z]+){0,2})\b'
        concepts = re.findall(concept_pattern, text)
        
        # Filter to get meaningful concepts (not common words)
        common_words = {'The', 'This', 'That', 'These', 'Those', 'Figure', 'Table', 'Page', 'Section'}
        unique_concepts = []
        seen = set()
        
        for concept in concepts:
            if (concept not in seen and 
                concept not in common_words and 
                len(concept) > 4 and
                text.count(concept) >= 2):  # Mentioned at least twice
                unique_concepts.append(concept)
                seen.add(concept)
        
        print(f"  Found {len(unique_concepts)} key concepts")
        
        # Create questions about concepts
        for concept in unique_concepts[:num - len(questions)]:
            # Find the most informative sentence about this concept
            concept_sentences = [s for s in sentences if concept in s]
            
            if concept_sentences:
                # Use the longest sentence as it likely has more context
                best_sentence = max(concept_sentences, key=len)
                
                # Truncate if too long
                if len(best_sentence) > 150:
                    best_sentence = best_sentence[:147] + "..."
                
                # Create wrong answers from other sentences
                other_content = [s[:100] for s in sentences if concept not in s][:3]
                
                wrong_answers = other_content if other_content else [
                    f"{concept} is not mentioned in the document",
                    f"{concept} is only briefly referenced",
                    f"{concept} appears in the appendix"
                ]
                
                options = [best_sentence] + wrong_answers[:3]
                random.shuffle(options)
                correct_idx = options.index(best_sentence)
                
                questions.append({
                    "question": f"What does the document explain about {concept}?",
                    "options": options,
                    "correct_answer": correct_idx,
                    "explanation": f"The document describes: {best_sentence}"
                })
                
                print(f"  ✓ Q: About {concept}")
        
        # If still need more questions, create comparison questions
        while len(questions) < num:
            if len(sentences) >= 4:
                # Pick a random informative sentence
                sent = random.choice([s for s in sentences if len(s) > 80])
                
                # Extract first part as question basis
                parts = sent.split(',')
                if len(parts) > 1:
                    main_part = parts[0]
                    
                    # Create options
                    other_parts = [s.split(',')[0] for s in random.sample(sentences, min(3, len(sentences)))]
                    
                    options = [sent[:120]] + [p[:120] for p in other_parts[:3]]
                    random.shuffle(options)
                    correct_idx = options.index(sent[:120])
                    
                    questions.append({
                        "question": f"Which statement from the document is accurate regarding {main_part}?",
                        "options": options,
                        "correct_answer": correct_idx,
                        "explanation": f"The document states: {sent}"
                    })
                else:
                    break
            else:
                break
        
        print(f"✅ Generated {len(questions)} intelligent quiz questions")
        return questions[:num]

# ================== Document Processor ==================
class DocumentProcessor:
    @staticmethod
    async def extract_text(file: UploadFile) -> str:
        content = await file.read()
        filename = file.filename.lower()
        
        try:
            if filename.endswith('.pdf') and HAS_PDF:
                pdf_reader = PyPDF2.PdfReader(io.BytesIO(content))
                text = ""
                for page in pdf_reader.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text += page_text + "\n"
                return text if text.strip() else "No text extracted"
            
            elif filename.endswith('.docx') and HAS_DOCX:
                doc = Document(io.BytesIO(content))
                return '\n\n'.join([p.text for p in doc.paragraphs if p.text.strip()])
            
            elif filename.endswith(('.txt', '.md')):
                return content.decode('utf-8', errors='ignore')
            
            return content.decode('utf-8', errors='ignore')
        except Exception as e:
            raise HTTPException(400, f"Text extraction failed: {str(e)}")

# ================== Presentation Generator ==================
class PresentationGenerator:
    @staticmethod
    def create_presentation(title: str, content: Dict) -> bytes:
        if not HAS_PPTX:
            raise Exception("Install python-pptx: pip install python-pptx")
        
        prs = Presentation()
        
        # Title slide
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = title
        if len(slide.placeholders) > 1:
            slide.placeholders[1].text = f"AI Analysis\n{datetime.now().strftime('%B %d, %Y')}"
        
        # Summary
        if content.get('summary'):
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = "Summary"
            tf = slide.placeholders[1].text_frame
            summary = content['summary'][:800]
            tf.text = summary
        
        # Key Points
        if content.get('key_points'):
            points = content['key_points']
            for i in range(0, len(points), 6):
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                slide.shapes.title.text = "Key Points" + (f" ({i//6+1})" if len(points) > 6 else "")
                tf = slide.placeholders[1].text_frame
                tf.clear()
                
                for j, point in enumerate(points[i:i+6]):
                    if len(point) > 150:
                        point = point[:147] + "..."
                    p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                    p.text = f"• {point}"
                    p.font.size = Pt(16)
        
        # Quiz
        if content.get('quiz'):
            for i, q in enumerate(content['quiz'], 1):
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                slide.shapes.title.text = f"Question {i}"
                tf = slide.placeholders[1].text_frame
                tf.clear()
                
                p = tf.paragraphs[0]
                p.text = q['question']
                p.font.bold = True
                p.font.size = Pt(18)
                
                for j, opt in enumerate(q['options'][:4]):
                    p = tf.add_paragraph()
                    p.text = f"{chr(65+j)}. {opt}"
                    p.font.size = Pt(16)
                    if j == q.get('correct_answer', 0):
                        p.font.bold = True
                        p.font.color.rgb = RGBColor(0, 128, 0)
        
        # End slide
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "Thank You"
        
        output = io.BytesIO()
        prs.save(output)
        output.seek(0)
        return output.read()

# ================== Pipeline Executor ==================
class PipelineExecutor:
    def __init__(self):
        self.ai_service = RealAIService()
        self.presentation_gen = PresentationGenerator()
    
    async def execute(self, prompt: str, document_text: str) -> Dict:
        if len(document_text.strip()) < 50:
            raise HTTPException(400, "Document too short")
        
        import time
        start_time = time.time()

        pipeline_id = f"pipeline_{uuid.uuid4().hex[:12]}"
        results = {
            'pipeline_id': pipeline_id,
            'status': 'processing',
            'tasks': [],
            'outputs': {},
            'created_at': datetime.now().isoformat()
        }
        
        prompt_lower = prompt.lower()
        
        # Summarization
        if 'summar' in prompt_lower or not any(w in prompt_lower for w in ['quiz', 'question', 'presentation', 'key', 'point']):
            print("📝 Summarizing...")
            try:
                summary = await self.ai_service.summarize(document_text)
                results['outputs']['summary'] = summary
                results['tasks'].append({
                    'name': 'Summarization',
                    'status': 'completed',
                    'output_preview': 'Summary complete'
                })
            except Exception as e:
                print(f"Summary error: {e}")
        
        # Key Points
        if any(w in prompt_lower for w in ['key', 'point', 'extract', 'main', 'important']) or not any(w in prompt_lower for w in ['quiz', 'question', 'presentation']):
            print("🔍 Extracting key points...")
            try:
                key_points = await self.ai_service.extract_key_points(document_text)
                results['outputs']['key_points'] = key_points
                results['tasks'].append({
                    'name': 'Key Points Extraction',
                    'status': 'completed',
                    'output_preview': f'{len(key_points)} key points extracted',
                    'result': key_points  # Add the actual result here too
                })
                
                # Print them so you can see
                print("\n📋 KEY POINTS EXTRACTED:")
                for i, point in enumerate(key_points, 1):
                    print(f"   {i}. {point}")
                print()
                
            except Exception as e:
                print(f"Key points error: {e}")
                import traceback
                traceback.print_exc()
        
        # Quiz
        if any(w in prompt_lower for w in ['quiz', 'question', 'test']):
            print("❓ Generating quiz...")
            try:
                nums = re.findall(r'\d+', prompt)
                num_q = min(int(nums[0]), 10) if nums else 5
                quiz = await self.ai_service.generate_quiz(document_text, num_q)
                results['outputs']['quiz'] = quiz
                results['tasks'].append({
                    'name': 'Quiz Generation',
                    'status': 'completed',
                    'output_preview': f'{len(quiz)} questions'
                })
            except Exception as e:
                print(f"Quiz error: {e}")
        
        # Presentation
        if any(w in prompt_lower for w in ['presentation', 'slide', 'ppt', 'powerpoint']):
            print("🎯 Creating presentation...")
            try:
                ppt_bytes = self.presentation_gen.create_presentation(
                    "Document Analysis",
                    {
                        'summary': results['outputs'].get('summary', ''),
                        'key_points': results['outputs'].get('key_points', []),
                        'quiz': results['outputs'].get('quiz', [])
                    }
                )
                results['status'] = 'completed'
                execution_time = time.time() - start_time
                results['execution_time'] = round(execution_time, 2)
                results['execution_time_formatted'] = f"{execution_time:.1f}s"
    
                pipelines[pipeline_id] = results
    
                return results
                
                filename = f"presentation_{pipeline_id}.pptx"
                filepath = Path("/tmp") / filename
                filepath.parent.mkdir(exist_ok=True)
                
                with open(filepath, 'wb') as f:
                    f.write(ppt_bytes)
                
                results['outputs']['presentation'] = {
                    'filename': filename,
                    'size': len(ppt_bytes),
                    'download_url': f"/download/{filename}"
                }
                results['tasks'].append({
                    'name': 'Presentation',
                    'status': 'completed',
                    'output_preview': 'PowerPoint ready'
                })
            except Exception as e:
                print(f"PPT error: {e}")
        
        results['status'] = 'completed'
        pipelines[pipeline_id] = results
        return results

pipeline_executor = PipelineExecutor()

# ================== API Endpoints ==================
@app.get("/")
async def root():
    return {
        "name": "Prompt-to-Pipeline API",
        "version": "4.0.0",
        "ai_method": pipeline_executor.ai_service.method,
        "endpoints": {
            "upload": "/upload",
            "execute": "/pipeline/execute",
            "download": "/download/{filename}"
        }
    }

@app.post("/upload")
async def upload_document(file: UploadFile = File(...)):
    content = await file.read()
    if len(content) > 10 * 1024 * 1024:
        raise HTTPException(413, "File too large (max 10MB)")
    
    file.file.seek(0)
    
    try:
        processor = DocumentProcessor()
        text = await processor.extract_text(file)
        
        if len(text.strip()) < 50:
            raise HTTPException(400, "Insufficient text extracted")
        
        doc_id = f"doc_{uuid.uuid4().hex[:12]}"
        uploaded_documents[doc_id] = {
            'filename': file.filename,
            'text': text,
            'size': len(content),
            'upload_time': datetime.now().isoformat(),
            'char_count': len(text),
            'word_count': len(text.split())
        }
        
        return {
            'document_id': doc_id,
            'filename': file.filename,
            'text_length': len(text),
            'word_count': len(text.split()),
            'preview': text[:500] + '...' if len(text) > 500 else text
        }
    except Exception as e:
        raise HTTPException(500, f"Processing error: {str(e)}")

@app.post("/pipeline/execute")
async def execute_pipeline_endpoint(request: Dict[str, Any]):
    prompt = request.get('prompt', 'Summarize and analyze')
    document_id = request.get('document_id')
    document_text = request.get('document_text', '')
    
    if document_id and document_id in uploaded_documents:
        document_text = uploaded_documents[document_id]['text']
    
    if not document_text:
        raise HTTPException(400, "No document provided")
    
    try:
        results = await pipeline_executor.execute(prompt, document_text)
        
        # DEBUG: Print what we're returning
        print("\n" + "="*60)
        print("PIPELINE RESULTS:")
        print("="*60)
        if 'summary' in results.get('outputs', {}):
            print(f"✓ Summary: {len(results['outputs']['summary'])} chars")
            print(f"  Preview: {results['outputs']['summary'][:200]}...")
        if 'key_points' in results.get('outputs', {}):
            print(f"✓ Key Points: {len(results['outputs']['key_points'])} points")
            for i, point in enumerate(results['outputs']['key_points'][:3], 1):
                print(f"  {i}. {point[:100]}...")
        if 'quiz' in results.get('outputs', {}):
            print(f"✓ Quiz: {len(results['outputs']['quiz'])} questions")
            for i, q in enumerate(results['outputs']['quiz'][:2], 1):
                print(f"  Q{i}: {q['question'][:80]}...")
                print(f"      Options: {q['options']}")
        print("="*60 + "\n")
        
        return results
    except Exception as e:
        print(f"ERROR in pipeline execution: {e}")
        import traceback
        traceback.print_exc()
        raise HTTPException(500, f"Execution failed: {str(e)}")

@app.get("/download/{filename}")
async def download_file(filename: str):
    safe_filename = filename.replace('..', '').replace('/', '')
    filepath = Path("/tmp") / safe_filename
    
    if not filepath.exists():
        raise HTTPException(404, f"File not found: {filename}")
    
    return FileResponse(
        path=str(filepath),
        media_type='application/vnd.openxmlformats-officedocument.presentationml.presentation',
        filename=safe_filename,
        headers={'Content-Disposition': f'attachment; filename="{safe_filename}"'}
    )

@app.get("/health")
async def health_check():
    return {
        'status': 'healthy',
        'ai_method': pipeline_executor.ai_service.method,
        'documents': len(uploaded_documents),
        'pipelines': len(pipelines)
    }

@app.get("/test/quiz/{document_id}")
async def test_quiz_generation(document_id: str, num_questions: int = 3):
    """Test endpoint to see quiz generation in detail"""
    if document_id not in uploaded_documents:
        raise HTTPException(404, "Document not found")
    
    text = uploaded_documents[document_id]['text']
    quiz = await pipeline_executor.ai_service.generate_quiz(text, num_questions)
    
    return {
        'method': pipeline_executor.ai_service.method,
        'num_questions': len(quiz),
        'quiz': quiz
    }

@app.get("/test/keypoints/{document_id}")
async def test_keypoints(document_id: str):
    """Test endpoint to see key points extraction"""
    if document_id not in uploaded_documents:
        raise HTTPException(404, "Document not found")
    
    text = uploaded_documents[document_id]['text']
    points = await pipeline_executor.ai_service.extract_key_points(text)
    
    return {
        'method': pipeline_executor.ai_service.method,
        'num_points': len(points),
        'key_points': points
    }

if __name__ == "__main__":
    import uvicorn
    print("=" * 70)
    print("🚀 PROMPT-TO-PIPELINE - GEMINI 2.5 ENHANCED")
    print("=" * 70)
    print(f"🤖 AI Method: {pipeline_executor.ai_service.method.upper()}")
    
    if pipeline_executor.ai_service.method == "gemini":
        if hasattr(pipeline_executor.ai_service, 'use_rest_api') and pipeline_executor.ai_service.use_rest_api:
            model = getattr(pipeline_executor.ai_service, 'gemini_model', 'gemini-2.5-flash')
            print(f"✅ Using Gemini REST API: {model}")
            print("   Summary: ✅ AI-powered")
            print("   Key Points: ✅ AI-powered")
            print("   Quiz: ✅ AI-powered")
        else:
            print("⚠️ Gemini SDK (might not work)")
    else:
        print("⚠️ Using Enhanced Fallback (no AI)")
        print("   Summary: Rule-based extraction")
        print("   Key Points: Keyword extraction")
        print("   Quiz: Context-based generation")
    
    print("=" * 70)
    print("🌐 http://localhost:8000")
    print("📚 API Docs: http://localhost:8000/docs")
    print("=" * 70)
    
    uvicorn.run(app, host="0.0.0.0", port=8000, reload=True)